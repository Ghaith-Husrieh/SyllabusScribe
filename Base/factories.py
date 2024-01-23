import random
from io import BytesIO
from uuid import uuid4

from django.contrib.auth.hashers import make_password
from django.core.files.base import ContentFile
from factory import (Factory, Faker, RelatedFactory, RelatedFactoryList,
                     SubFactory, lazy_attribute)
from factory.django import DjangoModelFactory
from factory.faker import faker
from pptx import Presentation

from .models import (LessonContext, LessonHandout, LessonPlan,
                     LessonPresentation, LessonQuiz, QuizQA, Subject, Unit,
                     User)

FAKE = faker.Faker()


class PowerPointFactory(Factory):
    class Meta:
        model = Presentation

    @classmethod
    def _create(cls, model_class, *args, **kwargs):
        presentation = model_class()

        cls.create_title_slide(presentation)
        cls.create_content_slide(presentation)

        pptx_data = BytesIO()
        presentation.save(pptx_data)
        unique_identifier = str(uuid4())
        file_name = f"Fake_Presentation_{unique_identifier[:8]}.pptx"

        return ContentFile(pptx_data.getvalue(), name=file_name)

    @classmethod
    def create_title_slide(cls, presentation):
        title_slide_layout = presentation.slide_layouts[0]
        title_slide = presentation.slides.add_slide(title_slide_layout)
        title_box = title_slide.shapes.title
        title_box.text = FAKE.sentence(nb_words=6)
        sub_title = title_slide.placeholders[1]
        sub_title.text = "Created by SyllabusScribe"

    @classmethod
    def create_content_slide(cls, presentation):
        content_slide_layout = presentation.slide_layouts[1]
        content_slide = presentation.slides.add_slide(content_slide_layout)
        content_box = content_slide.shapes.title
        content_box.text = FAKE.paragraph(nb_sentences=10)


class UserFactory(DjangoModelFactory):
    class Meta:
        model = User

    first_name = Faker('first_name')
    last_name = Faker('last_name')
    username = Faker('user_name')
    email = Faker('email')

    @lazy_attribute
    def password(self):
        # raw_password = FAKE.password()
        raw_password = 'testing1234'  # for easier accessibility to generated users
        hashed_password = make_password(raw_password)
        return hashed_password


class SubjectFactory(DjangoModelFactory):
    class Meta:
        model = Subject

    name = Faker('sentence', nb_words=6)
    user = SubFactory(UserFactory)


class LessonPlanFactory(DjangoModelFactory):
    class Meta:
        model = LessonPlan

    topic = Faker('sentence', nb_words=6)
    grade_level = Faker('word')
    user = SubFactory(UserFactory)

    @lazy_attribute
    def content(self):
        generated_content = f"{FAKE.paragraph(nb_sentences=15)}"
        for _ in range(0, 4):
            generated_content += f"\n\n{FAKE.paragraph(nb_sentences=15)}"
        return generated_content


class LessonContextFactory(DjangoModelFactory):
    class Meta:
        model = LessonContext

    topic = Faker('sentence', nb_words=6)
    grade_level = Faker('word')
    user = SubFactory(UserFactory)

    @lazy_attribute
    def content(self):
        generated_content = f"{FAKE.paragraph(nb_sentences=15)}"
        for _ in range(0, 4):
            generated_content += f"\n\n{FAKE.paragraph(nb_sentences=15)}"
        return generated_content


class LessonPresentationFactory(DjangoModelFactory):
    class Meta:
        model = LessonPresentation

    topic = Faker('sentence', nb_words=6)
    grade_level = Faker('word')
    user = SubFactory(UserFactory)
    generated_file = SubFactory(PowerPointFactory)


class LessonHandoutFactory(DjangoModelFactory):
    class Meta:
        model = LessonHandout

    topic = Faker('sentence', nb_words=6)
    grade_level = Faker('word')
    user = SubFactory(UserFactory)

    @lazy_attribute
    def content(self):
        generated_content = f"{FAKE.paragraph(nb_sentences=15)}"
        for _ in range(0, 4):
            generated_content += f"\n\n{FAKE.paragraph(nb_sentences=15)}"
        return generated_content


class _QuizQAFactory(DjangoModelFactory):
    class Meta:
        model = QuizQA

    @lazy_attribute
    def question(self):
        generated_question = FAKE.sentence(nb_words=6)
        for i in ['a', 'b', 'c', 'd']:
            generated_question += f"\n{i}) {FAKE.sentence(nb_words=6)}"
        return generated_question

    @lazy_attribute
    def answer(self):
        return f"{random.choice(['a', 'b', 'c', 'd'])}) {FAKE.sentence(nb_words=6)}"


class LessonQuizFactory(DjangoModelFactory):
    class Meta:
        model = LessonQuiz

    topic = Faker('sentence', nb_words=6)
    grade_level = Faker('word')
    user = SubFactory(UserFactory)

    quiz_qas = RelatedFactoryList(_QuizQAFactory, 'lesson_quiz', size=4)


class UnitFactory(DjangoModelFactory):
    class Meta:
        model = Unit

    name = Faker('sentence', nb_words=6)
    description = Faker('paragraph', nb_sentences=5)
    subject = SubFactory(SubjectFactory)

    @lazy_attribute
    def lesson_plan(self):
        return LessonPlanFactory(user=None)

    @lazy_attribute
    def lesson_context(self):
        return LessonContextFactory(user=None)

    @lazy_attribute
    def lesson_presentation(self):
        return LessonPresentationFactory(user=None)

    @lazy_attribute
    def lesson_handout(self):
        return LessonHandoutFactory(user=None)

    @lazy_attribute
    def lesson_quiz(self):
        return LessonQuizFactory(user=None)
