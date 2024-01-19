import re
from io import BytesIO
from uuid import uuid4

from django.core.files.base import ContentFile
from guidance import assistant, gen, system, user
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from decorators.log_decorators import log_function

from .prompt_constants import PRESENTATION_SYSTEM_PROMPT

TITLE_FONT_SIZE = Pt(32)
CONTENT_FONT_SIZE = Pt(16)


@log_function(log_result=False)
def generate_presentation_titles(model, num_slides, topic, grade_level):
    prompt = f"Generate {num_slides} short slide titles for the topic '{topic}' for {grade_level} students."
    with system():
        model += PRESENTATION_SYSTEM_PROMPT
    with user():
        model += prompt
    with assistant():
        model += f"here are {num_slides} potential slide titles for a presentation on '{topic}' for {grade_level} students:\n" + \
            gen("output", max_tokens=2048, temperature=0.0, top_p=0.1, stop="\n\n")
    generated_titles = model["output"].split("\n")
    for index, title in enumerate(generated_titles):
        generated_titles[index] = re.sub(r'[\d.\"]', '', title)
    return generated_titles


@log_function(log_result=False)
def generate_title_content(model, slide_title):
    prompt = f"Generate content for the slide: {slide_title}. The content must be in medium worded paragraphs. Only return 1 paragraph."
    with system():
        model += PRESENTATION_SYSTEM_PROMPT
    with user():
        model += prompt
    with assistant():
        model += f"here is a paragraph of content for a slide on '{slide_title}':\n" + \
            gen("output", max_tokens=320, temperature=0.0, top_p=0.1)
    generated_content = re.sub(r'[\"]', '', model["output"])
    return generated_content


@log_function(log_result=False)
def generate_presentation(topic, presentation_titles, presentation_contents):
    powerpoint = Presentation()

    title_slide_layout = powerpoint.slide_layouts[0]
    content_slide_layout = powerpoint.slide_layouts[1]

    background_color = RGBColor(255, 255, 255)

    title_slide = powerpoint.slides.add_slide(title_slide_layout)

    title = title_slide.shapes.title
    title.text = topic
    title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title.text_frame.paragraphs[0].font.bold = True

    content = title_slide.placeholders[1]
    content.text = "Created by SyllabusScribe"

    background = title_slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = background_color

    for slide_title, slide_content in zip(presentation_titles, presentation_contents):
        slide = powerpoint.slides.add_slide(content_slide_layout)

        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = background_color

        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        title.text_frame.paragraphs[0].font.bold = True

        content = slide.placeholders[1]
        content.text = slide_content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = CONTENT_FONT_SIZE

    pptx_data = BytesIO()
    powerpoint.save(pptx_data)
    unique_identifier = str(uuid4())
    file_name = f"{topic}_{unique_identifier[:8]}.pptx"

    return ContentFile(pptx_data.getvalue(), name=file_name)
